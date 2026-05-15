"""Generate a Sales Quarterly Update PowerPoint presentation from data.

Uses python-pptx to build a clean executive-style deck with text + plotly-rendered
PNG charts.
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from agent.visualizations import (
    pipeline_funnel,
    revenue_by_industry_donut,
    revenue_by_quarter_line,
    top_accounts_bar,
)

REPO_ROOT = Path(__file__).resolve().parent.parent
OUTPUTS = REPO_ROOT / "outputs"


PRIMARY = RGBColor(0x26, 0x4E, 0x86)
ACCENT = RGBColor(0xE2, 0x6A, 0x2C)
LIGHT_GREY = RGBColor(0x55, 0x55, 0x55)
BG = RGBColor(0xF7, 0xF7, 0xF8)


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Title
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12.5), Inches(1.6))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(46)
    run.font.bold = True
    run.font.color.rgb = PRIMARY
    # Subtitle
    sx = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12.5), Inches(0.8))
    sp = sx.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    sub = sp.add_run()
    sub.text = subtitle
    sub.font.size = Pt(20)
    sub.font.color.rgb = LIGHT_GREY


def _add_metric_slide(prs: Presentation, title: str, metrics: list[tuple[str, str, str]]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_title_bar(slide, title)
    n = len(metrics)
    col_w = 12.0 / max(n, 1)
    for i, (label, value, sub) in enumerate(metrics):
        left = Inches(0.6 + i * col_w)
        box = slide.shapes.add_textbox(left, Inches(2.0), Inches(col_w), Inches(3.0))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = label.upper()
        r1.font.size = Pt(12); r1.font.color.rgb = LIGHT_GREY; r1.font.bold = True
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = value
        r2.font.size = Pt(40); r2.font.bold = True; r2.font.color.rgb = PRIMARY
        if sub:
            p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run(); r3.text = sub
            r3.font.size = Pt(14); r3.font.color.rgb = LIGHT_GREY


def _add_title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(12.7), Inches(0.8))
    p = bar.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = PRIMARY


def _add_chart_slide(prs: Presentation, title: str, png_path: Path | None, caption: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_title_bar(slide, title)
    if png_path and Path(png_path).exists():
        slide.shapes.add_picture(str(png_path), Inches(0.5), Inches(1.3), width=Inches(12.3), height=Inches(5.5))
    if caption:
        cx = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.4))
        p = cx.text_frame.paragraphs[0]
        r = p.add_run(); r.text = caption
        r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = LIGHT_GREY


def _add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], note: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_title_bar(slide, title)
    nrows = len(rows) + 1
    ncols = len(headers)
    table_shape = slide.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5))
    table = table_shape.table
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    for ri, row in enumerate(rows):
        for ci, v in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11); r.font.color.rgb = LIGHT_GREY
    if note:
        cx = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.0), Inches(0.4))
        p = cx.text_frame.paragraphs[0]
        r = p.add_run(); r.text = note
        r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = LIGHT_GREY


def build_quarterly_update_deck(payload: dict[str, Any], out_path: Path | None = None) -> Path:
    """Build a Sales Quarterly Update from the output of analytics_quarterly_executive_update.

    `payload` is the dict returned by that tool (see analytics/tools.py).
    """
    out_path = out_path or (OUTPUTS / f"sales_quarterly_update_{payload['period']}.pptx")
    OUTPUTS.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Title slide
    _add_title_slide(prs, f"Sales Quarterly Update — {payload['period']}",
                     f"Generated {datetime.utcnow().strftime('%Y-%m-%d')} by MCP-powered enterprise AI agent")

    # 2. Headline metrics
    rev = payload["totals"]["revenue"]
    prior = payload["totals"].get("prior_year_revenue")
    yoy = payload["totals"].get("yoy_growth_pct")
    metrics = [
        ("Total Revenue", f"${rev:,.0f}", payload["period"]),
        ("Prior Year", f"${prior:,.0f}" if prior else "—", payload.get("prior_period", "")),
        ("YoY Growth", f"{yoy:+.1f}%" if yoy is not None else "—", ""),
        ("Open Pipeline",
         f"${sum(p['amount'] for p in payload['pipeline_by_stage']):,.0f}",
         f"{sum(p['n'] for p in payload['pipeline_by_stage'])} opps"),
    ]
    _add_metric_slide(prs, "Headline Metrics", metrics)

    # 3. Revenue by Industry (donut)
    if payload["revenue_by_industry"]:
        png, _ = revenue_by_industry_donut(payload["revenue_by_industry"], payload["period"],
                                            OUTPUTS / "_chart_industry")
        _add_chart_slide(prs, "Revenue by Industry", png, f"Period: {payload['period']}")

    # 4. Revenue by Product Family (table)
    if payload["revenue_by_product_family"]:
        rows = [[r["product_family"], f"${r['revenue']:,.0f}"] for r in payload["revenue_by_product_family"]]
        _add_table_slide(prs, "Revenue by Product Family", ["Product Family", "Revenue"], rows)

    # 5. Top 10 Customers
    if payload["top_customers"]:
        rows = [[i + 1, r["customer_name"], f"${r['revenue']:,.0f}"]
                for i, r in enumerate(payload["top_customers"])]
        _add_table_slide(prs, "Top Customers — Current Quarter", ["#", "Customer", "Revenue"], rows)

    # 6. Open Pipeline by Stage
    if payload["pipeline_by_stage"]:
        png, _ = pipeline_funnel(payload["pipeline_by_stage"], OUTPUTS / "_chart_pipeline")
        _add_chart_slide(prs, "Open Pipeline by Stage", png,
                         f"Total: ${sum(p['amount'] for p in payload['pipeline_by_stage']):,.0f}")

    # 7. Closed Won/Lost
    if payload["closed_won_lost"]:
        rows = [[r["stage"], r["n"], f"${r['amount']:,.0f}"] for r in payload["closed_won_lost"]]
        _add_table_slide(prs, "Closed Won / Lost — Current Quarter", ["Outcome", "Count", "Amount"], rows)

    # 8. Reliability concerns
    if payload["top_reliability_concerns"]:
        rows = [[i + 1, r["external_product_id"], r["rmas"], f"${r['cost']:,.0f}"]
                for i, r in enumerate(payload["top_reliability_concerns"])]
        _add_table_slide(prs, "Top Reliability Concerns This Quarter",
                         ["#", "Product ID", "RMAs", "Replacement Cost"], rows,
                         note="From QA system — products with the most customer returns this quarter.")

    # 9. Closing slide
    _add_metric_slide(prs, "Key Takeaways", [
        ("Revenue", f"${rev:,.0f}", payload["period"]),
        ("YoY", f"{yoy:+.1f}%" if yoy is not None else "—", "vs. prior year"),
        ("Top Customer", payload["top_customers"][0]["customer_name"] if payload["top_customers"] else "—",
         f"${payload['top_customers'][0]['revenue']:,.0f}" if payload["top_customers"] else ""),
        ("Pipeline Health", f"{len(payload['pipeline_by_stage'])} stages",
         f"${sum(p['amount'] for p in payload['pipeline_by_stage']):,.0f}"),
    ])

    prs.save(out_path)
    return out_path
